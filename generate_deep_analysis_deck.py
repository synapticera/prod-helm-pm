from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Create a presentation object
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # Define Brand Colors
    DARK_NAVY = RGBColor(7, 37, 61)      # #07253D
    CYAN = RGBColor(1, 169, 219)          # #01A9DB
    ORANGE = RGBColor(251, 65, 0)         # #FB4100
    WHITE = RGBColor(255, 255, 255)       # #FFFFFF
    GRAY = RGBColor(68, 84, 106)          # #44546A
    LIGHT_GRAY = RGBColor(200, 200, 200)
    GREEN = RGBColor(46, 204, 113)

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_NAVY

    def add_title(slide, text, font_size=36):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = text
        p = title_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = CYAN

    def add_bullet_slide(title, bullets, sub_bullets=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_background(slide)
        add_title(slide, title)
        
        body_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
        tf = body_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(24)
            p.font.color.rgb = WHITE
            p.space_before = Pt(15)
            
            if sub_bullets and i < len(sub_bullets) and sub_bullets[i]:
                for sub in sub_bullets[i]:
                    p_sub = tf.add_paragraph()
                    p_sub.text = "   - " + sub
                    p_sub.font.size = Pt(18)
                    p_sub.font.color.rgb = LIGHT_GRAY
                    p_sub.level = 1

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide1)
    
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategic Path Analysis"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = "Services + Product Business Decision Framework"
    sub.font.size = Pt(28)
    sub.font.color.rgb = CYAN
    sub.alignment = PP_ALIGN.CENTER
    sub.space_before = Pt(20)
    
    date_p = tf.add_paragraph()
    date_p.text = "January 9, 2026"
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = GRAY
    date_p.alignment = PP_ALIGN.CENTER
    date_p.space_before = Pt(40)

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide2)
    add_title(slide2, "Executive Summary")
    
    # Left: The Situation
    box_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = RGBColor(20, 50, 80)
    box_left.line.color.rgb = CYAN
    
    tf_l = box_left.text_frame
    tf_l.margin_left = Inches(0.2)
    tf_l.margin_top = Inches(0.2)
    
    p = tf_l.paragraphs[0]
    p.text = "The Situation"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ORANGE
    
    items = [
        "Pivotal decision point with two assets:",
        "1. Services: $4M Salesforce consulting",
        "2. Product (Helm): AI decision infra",
        "Critical Question: Belief in Helm's market timing & execution?"
    ]
    for item in items:
        p = tf_l.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # Right: Recommendation
    box_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.5), Inches(4.5))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = RGBColor(20, 50, 80)
    box_right.line.color.rgb = GREEN
    
    tf_r = box_right.text_frame
    tf_r.margin_left = Inches(0.2)
    tf_r.margin_top = Inches(0.2)
    
    p = tf_r.paragraphs[0]
    p.text = "Recommendation"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = GREEN
    
    rec_items = [
        "Path 4: Strategic Leverage",
        "Use services strategically to de-risk product development.",
        "Maintain optionality.",
        "Not full lean, not full liquidation.",
        "A nuanced, phased approach."
    ]
    for item in rec_items:
        p = tf_r.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # --- Slide 3: Market Context ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide3)
    add_title(slide3, "Market Context: Why Now?")

    # AI Market Table Visual
    y_start = Inches(1.5)
    
    # Header
    headers = ["Factor", "Data Point", "Implication"]
    header_x = [Inches(1), Inches(4), Inches(8)]
    for i, h in enumerate(headers):
        tb = slide3.shapes.add_textbox(header_x[i], y_start, Inches(3), Inches(0.5))
        tb.text_frame.text = h
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = CYAN
        tb.text_frame.paragraphs[0].font.size = Pt(18)

    data = [
        ("Agent adoption", "40% of enterprise apps by 2026", "Massive deployment wave"),
        ("Governance gap", "33% of AI deployments fail", "Mgmt infra is urgent need"),
        ("Spending deferral", "25% of AI spend moving to 2027", "CFOs demanding ROI proof"),
        ("Regulatory", "EU AI Act effective Aug 2026", "Compliance is mandatory"),
        ("VC alignment", "Foundation Cap $4.6T thesis", "Decision infra is THE opp")
    ]

    curr_y = y_start + Inches(0.6)
    for row in data:
        # Draw line
        line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), curr_y, Inches(11.3), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = GRAY
        
        curr_y += Inches(0.1)
        
        # Col 1
        tb1 = slide3.shapes.add_textbox(header_x[0], curr_y, Inches(2.8), Inches(0.5))
        tb1.text_frame.text = row[0]
        tb1.text_frame.paragraphs[0].font.color.rgb = WHITE
        tb1.text_frame.paragraphs[0].font.bold = True
        
        # Col 2
        tb2 = slide3.shapes.add_textbox(header_x[1], curr_y, Inches(3.8), Inches(0.5))
        tb2.text_frame.text = row[1]
        tb2.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
        
        # Col 3
        tb3 = slide3.shapes.add_textbox(header_x[2], curr_y, Inches(4), Inches(0.5))
        tb3.text_frame.text = row[2]
        tb3.text_frame.paragraphs[0].font.color.rgb = ORANGE
        
        curr_y += Inches(0.6)

    footer = slide3.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
    footer.text_frame.text = "Insight: Market window is 12-24 months. Speed is critical."
    footer.text_frame.paragraphs[0].font.color.rgb = CYAN
    footer.text_frame.paragraphs[0].font.size = Pt(20)
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slide 4: Path 1 - Lean ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide4)
    add_title(slide4, "Path 1: LEAN")
    
    sub = slide4.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
    sub.text_frame.text = "Keep services lean & cash-flow positive. Fund product from profits."
    sub.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    sub.text_frame.paragraphs[0].font.size = Pt(18)

    # Pros/Cons Split
    y_split = Inches(2)
    
    # Pros
    tb_pro = slide4.shapes.add_textbox(Inches(0.5), y_split, Inches(5.5), Inches(0.5))
    tb_pro.text_frame.text = "PROS"
    tb_pro.text_frame.paragraphs[0].font.color.rgb = GREEN
    tb_pro.text_frame.paragraphs[0].font.bold = True
    
    pros = ["Focus (single direction)", "Capital efficiency (no dilution)", "Runway to PMF", "Optionality"]
    for p_text in pros:
        tb = slide4.shapes.add_textbox(Inches(0.5), y_split + Inches(0.5), Inches(5.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "• " + p_text
        p.font.color.rgb = WHITE
        y_split += Inches(0.5)

    # Cons
    y_split = Inches(2) # Reset y
    tb_con = slide4.shapes.add_textbox(Inches(6.5), y_split, Inches(5.5), Inches(0.5))
    tb_con.text_frame.text = "CONS"
    tb_con.text_frame.paragraphs[0].font.color.rgb = ORANGE
    tb_con.text_frame.paragraphs[0].font.bold = True
    
    cons = ["Speed disadvantage vs funded competitors", "Team morale risk", "Single point of failure", "Misses 12-24mo market window"]
    for c_text in cons:
        tb = slide4.shapes.add_textbox(Inches(6.5), y_split + Inches(0.5), Inches(5.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = "• " + c_text
        p.font.color.rgb = WHITE
        y_split += Inches(0.5)

    # Assessment
    assess_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
    assess_box.fill.solid()
    assess_box.fill.fore_color.rgb = RGBColor(100, 30, 30) # Dark red
    assess_box.line.color.rgb = ORANGE
    
    atf = assess_box.text_frame
    atf.margin_left = Inches(0.2)
    p = atf.paragraphs[0]
    p.text = "ASSESSMENT: High Risk"
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = atf.add_paragraph()
    p2.text = "Underestimates urgency. You might perfect the product just as the window closes."
    p2.font.color.rgb = WHITE

    # --- Slide 5: Path 2 - Liquidate ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide5)
    add_title(slide5, "Path 2: LIQUIDATE")
    
    sub = slide5.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
    sub.text_frame.text = "Sell services (~$5M net). Fund product aggressively."
    sub.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    sub.text_frame.paragraphs[0].font.size = Pt(18)

    # Financial Model Box
    fin_box = slide5.shapes.add_textbox(Inches(8), Inches(2), Inches(4), Inches(2))
    fin_box.text_frame.text = "Potential Outcomes:\nLow: $5.5M Net\nMid: $7M Net\nHigh: $10M Net"
    fin_box.text_frame.paragraphs[0].font.color.rgb = CYAN
    
    # Pros/Cons
    y_split = Inches(2)
    tb_pro = slide5.shapes.add_textbox(Inches(0.5), y_split, Inches(3.5), Inches(0.5))
    tb_pro.text_frame.text = "PROS"
    tb_pro.text_frame.paragraphs[0].font.color.rgb = GREEN
    
    pros = ["Full Focus", "Capitalization", "Speed", "Clean Narrative"]
    y_curr = y_split + Inches(0.4)
    for p_text in pros:
        tb = slide5.shapes.add_textbox(Inches(0.5), y_curr, Inches(3.5), Inches(0.4))
        tb.text_frame.text = "• " + p_text
        tb.text_frame.paragraphs[0].font.color.rgb = WHITE
        y_curr += Inches(0.4)
        
    tb_con = slide5.shapes.add_textbox(Inches(4), y_split, Inches(3.5), Inches(0.5))
    tb_con.text_frame.text = "CONS"
    tb_con.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    cons = ["Transaction distraction (6-12mo)", "Cash timing (Earnouts)", "All-in risk", "Lose 'Lab' value"]
    y_curr = y_split + Inches(0.4)
    for c_text in cons:
        tb = slide5.shapes.add_textbox(Inches(4), y_curr, Inches(3.5), Inches(0.4))
        tb.text_frame.text = "• " + c_text
        tb.text_frame.paragraphs[0].font.color.rgb = WHITE
        y_curr += Inches(0.4)

    # Assessment
    assess_box = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.5))
    assess_box.fill.solid()
    assess_box.fill.fore_color.rgb = RGBColor(80, 60, 20) # Brownish
    assess_box.line.color.rgb = ORANGE
    
    atf = assess_box.text_frame
    atf.margin_left = Inches(0.2)
    p = atf.paragraphs[0]
    p.text = "ASSESSMENT: High Reward / High Execution Risk"
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = atf.add_paragraph()
    p2.text = "Theoretically optimal but practically risky unless buyer is ready NOW. M&A is a huge distraction."
    p2.font.color.rgb = WHITE

    # --- Slide 6: Path 3 - Growth Both ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide6)
    add_title(slide6, "Path 3: GROWTH BOTH")
    
    sub = slide6.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
    sub.text_frame.text = "Hire services leader. Raise VC for product. Grow both."
    sub.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    sub.text_frame.paragraphs[0].font.size = Pt(18)

    # Pros/Cons
    y_split = Inches(2)
    tb_pro = slide6.shapes.add_textbox(Inches(0.5), y_split, Inches(5), Inches(0.5))
    tb_pro.text_frame.text = "PROS"
    tb_pro.text_frame.paragraphs[0].font.color.rgb = GREEN
    
    pros = ["Diversification", "Services as Lab", "Cash flow support"]
    y_curr = y_split + Inches(0.4)
    for p_text in pros:
        tb = slide6.shapes.add_textbox(Inches(0.5), y_curr, Inches(5), Inches(0.4))
        tb.text_frame.text = "• " + p_text
        tb.text_frame.paragraphs[0].font.color.rgb = WHITE
        y_curr += Inches(0.4)
        
    tb_con = slide6.shapes.add_textbox(Inches(6.5), y_split, Inches(5), Inches(0.5))
    tb_con.text_frame.text = "CONS"
    tb_con.text_frame.paragraphs[0].font.color.rgb = ORANGE
    
    cons = ["Focus Dilution (Very High)", "Capital Intensive", "VCs dislike hybrid models", "Culture clash"]
    y_curr = y_split + Inches(0.4)
    for c_text in cons:
        tb = slide6.shapes.add_textbox(Inches(6.5), y_curr, Inches(5), Inches(0.4))
        tb.text_frame.text = "• " + c_text
        tb.text_frame.paragraphs[0].font.color.rgb = WHITE
        y_curr += Inches(0.4)

    # VC Insight
    vc_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(12.33), Inches(1.2))
    vc_box.fill.solid()
    vc_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
    vc_box.line.color.rgb = GRAY
    vctf = vc_box.text_frame
    vctf.margin_left = Inches(0.2)
    vctf.text = "VC Reality: Pure product companies get 30-40x multiples. Hybrid gets 10-15x. Services revenue drags down valuation."
    vctf.paragraphs[0].font.color.rgb = LIGHT_GRAY

    # Assessment
    assess_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.2))
    assess_box.fill.solid()
    assess_box.fill.fore_color.rgb = RGBColor(100, 30, 30)
    assess_box.line.color.rgb = ORANGE
    
    atf = assess_box.text_frame
    atf.margin_left = Inches(0.2)
    p = atf.paragraphs[0]
    p.text = "ASSESSMENT: Overcomplicated"
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = atf.add_paragraph()
    p2.text = "VCs will likely push you to spin off anyway. Synergy is often overstated."
    p2.font.color.rgb = WHITE

    # --- Slide 7: Recommendation - Strategic Leverage ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide7)
    add_title(slide7, "Path 4: STRATEGIC LEVERAGE (Recommended)")
    
    sub = slide7.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
    sub.text_frame.text = "Use services as a bridge. Extract value, fund product validation, maintain optionality."
    sub.text_frame.paragraphs[0].font.color.rgb = CYAN
    sub.text_frame.paragraphs[0].font.size = Pt(20)

    # Phase 1
    p1_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(5.8), Inches(3.5))
    p1_box.fill.solid()
    p1_box.fill.fore_color.rgb = RGBColor(20, 50, 80)
    p1_box.line.color.rgb = CYAN
    
    tf = p1_box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "Phase 1: Strategic Extraction (Months 1-12)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = CYAN
    
    steps1 = [
        "1. Optimize services for cash (Target 20% EBITDA)",
        "2. Use services as 'Product Lab' (Design partners)",
        "3. Soft-market the services business",
        "4. Fund initial product dev (MVP)"
    ]
    for s in steps1:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)

    # Phase 2
    p2_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2), Inches(5.8), Inches(3.5))
    p2_box.fill.solid()
    p2_box.fill.fore_color.rgb = RGBColor(20, 50, 80)
    p2_box.line.color.rgb = GREEN
    
    tf = p2_box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = "Phase 2: Decision Point (Month 12-18)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = GREEN
    
    steps2 = [
        "A: Product Winning -> Accelerate sale, raise Series A",
        "B: Product Promising -> Maintain hybrid, seed raise",
        "C: Product Struggling -> Pivot or double-down on services"
    ]
    for s in steps2:
        p = tf.add_paragraph()
        p.text = s
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10)
        
    # Why this works
    footer = slide7.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(1))
    footer.text_frame.text = "Why: Preserves optionality, uses time wisely, de-risks for VCs, maintains focus."
    footer.text_frame.paragraphs[0].font.color.rgb = ORANGE
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer.text_frame.paragraphs[0].font.size = Pt(20)

    # --- Slide 8: Financial Expected Value ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide8)
    add_title(slide8, "Expected Value Analysis (3-Year)")

    # Table
    table_y = Inches(2)
    table_h = Inches(4)
    table_w = Inches(10)
    table_x = Inches(1.5)
    
    rows = 5
    cols = 2
    shape = slide8.shapes.add_table(rows, cols, table_x, table_y, table_w, table_h)
    table = shape.table
    
    # Headers
    table.cell(0, 0).text = "Strategic Path"
    table.cell(0, 1).text = "Expected Value (Probability Weighted)"
    
    data = [
        ("Path 1: Lean", "$17.8M"),
        ("Path 2: Liquidate", "$30.5M"),
        ("Path 3: Growth Both", "$32.2M"),
        ("Path 4: Strategic Leverage", "$37.8M")
    ]
    
    for i, (path, val) in enumerate(data):
        table.cell(i+1, 0).text = path
        table.cell(i+1, 1).text = val
        
    # Style table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 60, 90) if r > 0 else RGBColor(20, 20, 20)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    # Highlight winner
    table.cell(4, 1).text_frame.paragraphs[0].font.color.rgb = GREEN
    table.cell(4, 1).text_frame.paragraphs[0].font.bold = True

    # --- Slide 9: Implementation Roadmap ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide9)
    add_title(slide9, "Implementation Roadmap")
    
    # Timeline
    timeline_y = Inches(2.5)
    timeline_h = Inches(0.5)
    
    # Arrow line
    arrow = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1), timeline_y, Inches(11.3), timeline_h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY
    
    # Phase markers
    phases = [
        ("Months 1-3\nFoundation", Inches(1.5)),
        ("Months 4-6\nValidation", Inches(4.5)),
        ("Months 7-12\nDecision Point", Inches(7.5)),
        ("Months 13-24\nExecution", Inches(10.5))
    ]
    
    for text, x in phases:
        tb = slide9.shapes.add_textbox(x - Inches(1), timeline_y - Inches(1), Inches(2), Inches(1))
        tb.text_frame.text = text
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = CYAN
        
        # Details below
        details_map = {
            "Months 1-3\nFoundation": ["Optimize Cash", "Soft-market Svcs", "Identify Design Partners"],
            "Months 4-6\nValidation": ["Deploy MVP", "Case Studies", "Investor Pipeline"],
            "Months 7-12\nDecision Point": ["Assess PMF", "Choose: Sell/Raise/Pivot"],
            "Months 13-24\nExecution": ["Scale Product", "Optimize Exit"]
        }
        
        details = details_map[text]
        curr_dy = timeline_y + Inches(0.8)
        for d in details:
            dtb = slide9.shapes.add_textbox(x - Inches(1.2), curr_dy, Inches(2.4), Inches(0.5))
            dtb.text_frame.text = "• " + d
            dtb.text_frame.paragraphs[0].font.size = Pt(12)
            dtb.text_frame.paragraphs[0].font.color.rgb = WHITE
            dtb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            curr_dy += Inches(0.4)

    # --- Slide 10: Conclusion ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide10)
    add_title(slide10, "Conclusion & Next Steps")
    
    # Main Box
    box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2), Inches(10.33), Inches(3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(20, 50, 80)
    box.line.color.rgb = GREEN
    
    tf = box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.4)
    
    p = tf.paragraphs[0]
    p.text = "Don't go lean (too slow). Don't rush to liquidate (too risky). Don't split focus."
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Use services as a strategic bridge to your product future."
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = CYAN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(30)

    # Questions
    q_box = slide10.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.33), Inches(1.5))
    q_box.text_frame.text = "Decisions needed: Cash flow targets? RS buyer status? Product conviction?"
    q_box.text_frame.paragraphs[0].font.color.rgb = ORANGE
    q_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


    # Save
    output_path = "/Users/wkarp/devlocal/prod-miner-pm/output/Strategy_Deep_Analysis.pptx"
    prs.save(output_path)
    print(f"✅ Presentation generated successfully: {output_path}")

if __name__ == "__main__":
    create_presentation()
