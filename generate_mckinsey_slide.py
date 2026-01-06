from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add a blank slide
blank_slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(blank_slide_layout)

# Define Synaptic brand colors
DARK_NAVY = RGBColor(7, 37, 61)      # #07253D
CYAN = RGBColor(1, 169, 219)          # #01A9DB
ORANGE = RGBColor(251, 65, 0)         # #FB4100
WHITE = RGBColor(255, 255, 255)       # #FFFFFF
GRAY = RGBColor(68, 84, 106)          # #44546A

# Add dark navy background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_NAVY

# Add title at the top
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "McKinsey Insight: AI Deployment vs. Impact Gap"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = CYAN

# Left side: Statistics section
left_x = Inches(0.5)
stat_y = Inches(1.2)
stat_width = Inches(3)
stat_height = Inches(1)

# Stat 1: 88% AI Usage
stat1_box = slide.shapes.add_textbox(left_x, stat_y, stat_width, stat_height)
stat1_frame = stat1_box.text_frame
stat1_frame.word_wrap = True
p1 = stat1_frame.paragraphs[0]
p1.text = "88%"
p1.font.size = Pt(48)
p1.font.bold = True
p1.font.color.rgb = CYAN
p1.alignment = PP_ALIGN.CENTER

p2 = stat1_frame.add_paragraph()
p2.text = "Enterprise AI Usage"
p2.font.size = Pt(12)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER

# Stat 2: 39% Impact
stat2_box = slide.shapes.add_textbox(left_x, stat_y + Inches(1.1), stat_width, stat_height)
stat2_frame = stat2_box.text_frame
stat2_frame.word_wrap = True
p3 = stat2_frame.paragraphs[0]
p3.text = "39%"
p3.font.size = Pt(48)
p3.font.bold = True
p3.font.color.rgb = ORANGE
p3.alignment = PP_ALIGN.CENTER

p4 = stat2_frame.add_paragraph()
p4.text = "Seeing Measurable Impact"
p4.font.size = Pt(12)
p4.font.color.rgb = WHITE
p4.alignment = PP_ALIGN.CENTER

# Stat 3: Gap
stat3_box = slide.shapes.add_textbox(left_x, stat_y + Inches(2.2), stat_width, stat_height)
stat3_frame = stat3_box.text_frame
stat3_frame.word_wrap = True
p5 = stat3_frame.paragraphs[0]
p5.text = "-49 pts"
p5.font.size = Pt(48)
p5.font.bold = True
p5.font.color.rgb = RGBColor(255, 100, 100)  # Red
p5.alignment = PP_ALIGN.CENTER

p6 = stat3_frame.add_paragraph()
p6.text = "The Business Impact Gap"
p6.font.size = Pt(12)
p6.font.color.rgb = WHITE
p6.alignment = PP_ALIGN.CENTER

# Right side: Impact box with key insight
right_x = Inches(4)
impact_width = Inches(5.5)
impact_box = slide.shapes.add_shape(
    1,  # Rectangle shape
    right_x, stat_y, impact_width, Inches(3.2)
)
impact_box.fill.solid()
impact_box.fill.fore_color.rgb = RGBColor(30, 60, 90)  # Darker navy for contrast
impact_box.line.color.rgb = CYAN
impact_box.line.width = Pt(2)

# Add text inside impact box
text_frame = impact_box.text_frame
text_frame.margin_left = Inches(0.3)
text_frame.margin_right = Inches(0.3)
text_frame.margin_top = Inches(0.3)
text_frame.margin_bottom = Inches(0.3)
text_frame.word_wrap = True

p_title = text_frame.paragraphs[0]
p_title.text = "The Root Cause:"
p_title.font.size = Pt(14)
p_title.font.bold = True
p_title.font.color.rgb = CYAN
p_title.level = 0

p_insight = text_frame.add_paragraph()
p_insight.text = "Workflow Redesign"
p_insight.font.size = Pt(24)
p_insight.font.bold = True
p_insight.font.color.rgb = ORANGE
p_insight.level = 0
p_insight.space_before = Pt(6)

p_detail = text_frame.add_paragraph()
p_detail.text = "McKinsey research shows workflow redesign has the biggest effect on organizations' ability to see EBIT impact from generative AI."
p_detail.font.size = Pt(11)
p_detail.font.color.rgb = WHITE
p_detail.level = 0
p_detail.space_before = Pt(10)
p_detail.space_after = Pt(10)

p_solution = text_frame.add_paragraph()
p_solution.text = "Helm solves this through our Mine/Design/Manage framework—automating workflow redesign so organizations capture real business value."
p_solution.font.size = Pt(11)
p_solution.font.color.rgb = CYAN
p_solution.level = 0

# Add footer
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(0.4))
footer_frame = footer_box.text_frame
footer_frame.text = "McKinsey, 2024 | The State of AI"
footer_frame.paragraphs[0].font.size = Pt(10)
footer_frame.paragraphs[0].font.color.rgb = GRAY

# Save the presentation
output_path = "/Users/wkarp/devlocal/prod-miner-pm/output/McKinsey_Helm_Validation.pptx"
prs.save(output_path)

print(f"✅ Presentation generated successfully: {output_path}")
