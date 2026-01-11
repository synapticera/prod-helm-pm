#!/usr/bin/env python3
"""
Create EVOLVE investor deck PowerPoint from InvestorBrief.md,
preserving the existing title slide and styling.

Updates slides 2-14 while keeping the original title slide (slide 1) intact.
Archives old versions before creating new presentation.

14 slides total:
1. Title (preserved from existing)
2. Overview
3. The Problem
4. The Solution
5. The Tesla Parallel (Comparison Table)
6. Market Opportunity
7. Competitive Advantage
8. Business Model
9. Go-to-Market
10. Product Roadmap (Overview)
11. Product Roadmap by Quarter 2026 (Matrix Table)
12. Financial Projections
13. Team & Capital
14. The Ask
"""

import os
import shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt

def extract_styling_from_example(example_path):
    """Extract color scheme and fonts from example presentation."""
    try:
        prs = Presentation(example_path)
        theme = {}

        if len(prs.slides) > 0:
            slide = prs.slides[0]
            if slide.background.fill.type:
                theme['background'] = 'from_example'

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                theme['font_name'] = run.font.name
                            if run.font.size:
                                theme['font_size'] = run.font.size
                            if run.font.color.type == 1:
                                theme['text_color'] = run.font.color.rgb
                            break
                        if 'font_name' in theme:
                            break
                if 'font_name' in theme:
                    break

        print(f"Extracted theme: {theme}")
        return theme
    except Exception as e:
        print(f"Error extracting styling: {e}")
        return {}

def create_title_slide(prs, theme):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title = slide.shapes.title
    title.text = "EVOLVE"

    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Transforming Transformation\n\nSeries Seed Extension\nSeeking: $2-3M\n\nwww.evolve.ai"

    return slide

def create_content_slide(prs, title_text, content_list, theme):
    """Create a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)

    body = slide.placeholders[1]
    original_left = body.left
    original_width = body.width
    original_height = body.height
    current_top_inches = body.top / 914400
    body.top = int((current_top_inches - 0.5) * 914400)
    body.left = original_left
    body.width = original_width
    body.height = original_height
    tf = body.text_frame
    tf.clear()

    for item in content_list:
        if isinstance(item, dict):
            p = tf.add_paragraph()
            p.text = item['text']
            p.level = item.get('level', 0)
            if p.level == 0:
                p.space_before = Pt(12)
                p.space_after = Pt(2)
            else:
                p.space_before = Pt(2)
                p.space_after = Pt(2)
            p.line_spacing = 1.0
            for run in p.runs:
                if p.level == 0:
                    run.font.size = Pt(16)
                elif p.level == 1:
                    run.font.size = Pt(14)
                else:
                    run.font.size = Pt(12)
        else:
            p = tf.add_paragraph()
            p.text = str(item)
            p.level = 0
            p.space_before = Pt(12)
            p.space_after = Pt(2)
            p.line_spacing = 1.0
            for run in p.runs:
                run.font.size = Pt(16)

    return slide

def create_table_slide(prs, title_text, headers, rows, theme):
    """Create a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title
    title.text = title_text
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)

    rows_count = len(rows) + 1
    cols_count = len(headers)

    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    for col in range(cols_count):
        table.columns[col].width = Inches(width.inches / cols_count)

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.space_before = Pt(2)
            paragraph.space_after = Pt(2)

    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.space_before = Pt(2)
                paragraph.space_after = Pt(2)

    return slide

def parse_investor_brief_content():
    """Parse investor brief content into 12-slide deck structure."""
    slides_data = []

    # Slide 1: Title
    slides_data.append({
        'type': 'title',
        'title': 'EVOLVE',
        'subtitle': 'Transforming Transformation\n\nSeries Seed Extension\nSeeking: $2-3M\n\nwww.evolve.ai'
    })

    # Slide 2: Outline
    slides_data.append({
        'type': 'content',
        'title': 'Overview',
        'content': [
            {'text': 'The Problem: Transformation is Broken', 'level': 0},
            {'text': 'The Solution: End-to-End Platform', 'level': 0},
            {'text': 'Market Opportunity: $127B TAM', 'level': 0},
            {'text': 'Competitive Advantage: Network Effects', 'level': 0},
            {'text': 'Business Model: Freemium to Enterprise', 'level': 0},
            {'text': 'Go-to-Market: Salesforce SMB Entry', 'level': 0},
            {'text': 'Product Roadmap: 18-Month Path', 'level': 0},
            {'text': 'Projections: $2M to $300M ARR', 'level': 0},
            {'text': 'Team & Capital: $2-3M Seed Extension', 'level': 0},
            {'text': 'The Ask: Building Category Leader', 'level': 0},
        ]
    })

    # Slide 3: The Problem
    slides_data.append({
        'type': 'content',
        'title': 'The Problem: Transformation is Broken',
        'content': [
            {'text': '95% of transformation initiatives fail', 'level': 0},
            {'text': 'Companies flying blind - no visibility into actual processes', 'level': 1},
            {'text': 'Can\'t prioritize by ROI or compare to peers', 'level': 1},
            {'text': 'No continuous monitoring or assurance', 'level': 1},
            {'text': 'Traditional Model Can\'t Keep Up:', 'level': 0},
            {'text': 'Consulting: $2M+, 6-12 months, no execution', 'level': 1},
            {'text': 'Process Mining: $100K-500K, enterprise-only', 'level': 1},
            {'text': 'RPA: Brittle, no discovery, breaks constantly', 'level': 1},
            {'text': 'Result: 150K+ SMBs completely locked out', 'level': 0},
        ]
    })

    # Slide 4: The Solution
    slides_data.append({
        'type': 'content',
        'title': 'The EVOLVE Solution: Complete Lifecycle',
        'content': [
            {'text': '1. MINE (FREE): Edge-first workflow discovery', 'level': 0},
            {'text': 'Privacy-preserving, instant value, no security concerns', 'level': 1},
            {'text': '2. DESIGN: AI-powered transformation roadmaps', 'level': 0},
            {'text': 'Informed by aggregate patterns from thousands of companies', 'level': 1},
            {'text': '3. DELIVER: Autonomous agent authoring', 'level': 0},
            {'text': 'Shadow mode validation, continuous learning', 'level': 1},
            {'text': '4. MANAGE: Continuous improvement', 'level': 0},
            {'text': 'Real-time peer benchmarking, service assurance', 'level': 1},
            {'text': 'Result: Self-improving intelligence system', 'level': 0},
        ]
    })

    # Slide 5: The Tesla Parallel
    slides_data.append({
        'type': 'table',
        'title': 'The Tesla Parallel: Democratizing Intelligence',
        'headers': ["Tesla's Innovation", "EVOLVE's Innovation", "Impact"],
        'rows': [
            [
                "Commodity cameras\nvs. expensive LIDAR",
                "Commodity LLMs\nvs. specialized consultants",
                "Lowers costs 10x+\nAccessible to millions of SMBs"
            ],
            [
                "Edge processing\nin vehicles",
                "Edge-first architecture\non user devices",
                "70-80% lower costs\nZero security friction"
            ],
            [
                "Fleet learning:\nMillions of cars\nimprove autopilot",
                "Aggregate learning:\nThousands of companies\nimprove platform",
                "Compounding intelligence\nConsulting can't replicate"
            ],
            [
                "Freemium drove\nmass adoption",
                "Free discovery tier\ndrives adoption",
                "100x more users\nWinner-takes-most dynamics"
            ]
        ]
    })

    # Slide 6: Market Opportunity
    slides_data.append({
        'type': 'content',
        'title': 'Market Opportunity: Dual-Track TAM',
        'content': [
            {'text': 'Salesforce Service Cloud Entry Wedge:', 'level': 0},
            {'text': '$5.4B TAM (1.2-1.5M service agents)', 'level': 1},
            {'text': '$1.08B SAM (20K target SMB customers)', 'level': 1},
            {'text': '$71M-$212M SOM by Year 3-5', 'level': 1},
            {'text': 'Broader Technology Spend TAM:', 'level': 0},
            {'text': '$127B by 2030 (AI Agents $68B, Process Mining $23B, RPA $36B)', 'level': 1},
            {'text': '$31B SAM (US/UK/EU, 100-10K employees)', 'level': 1},
            {'text': 'Labor Savings TAM:', 'level': 0},
            {'text': '29M knowledge workers × $3.4K revenue = $98.6B annually', 'level': 1},
        ]
    })

    # Slide 6: Competitive Advantage
    slides_data.append({
        'type': 'content',
        'title': 'Competitive Advantage: Unbeatable Moat',
        'content': [
            {'text': 'Network Effects (Winner-Takes-Most):', 'level': 0},
            {'text': 'More users → Better benchmarks → Higher conversion', 'level': 1},
            {'text': 'Peer data doesn\'t exist anywhere else', 'level': 1},
            {'text': 'Switching costs increase (lose peer insights)', 'level': 1},
            {'text': 'Business Model Moat:', 'level': 0},
            {'text': 'Freemium = 100x more users than competitors', 'level': 1},
            {'text': 'Edge architecture = 70-80% lower costs', 'level': 1},
            {'text': 'SMB-first = $20K deals work (vs. $100K+ for them)', 'level': 1},
            {'text': 'Independence Moat: Platform-agnostic vs. vendor lock-in', 'level': 0},
        ]
    })

    # Slide 7: Business Model
    slides_data.append({
        'type': 'content',
        'title': 'Business Model: Build Fleet, Monetize Value',
        'content': [
            {'text': 'FREE: Build User Base', 'level': 0},
            {'text': 'Full workflow discovery, agent opportunities', 'level': 1},
            {'text': 'Goal: 100K+ users, aggregate intelligence', 'level': 1},
            {'text': 'PAID TIER 1: $500-$50K/mo', 'level': 0},
            {'text': 'Transformation execution, agent deployment', 'level': 1},
            {'text': 'Target: 30-40% conversion, SMB to mid-market', 'level': 1},
            {'text': 'PAID TIER 2: $50K-$500K+/mo', 'level': 0},
            {'text': 'Enterprise platform, continuous improvement', 'level': 1},
            {'text': 'Unit Economics: 229:1 LTV:CAC, 85% gross margin', 'level': 0},
        ]
    })

    # Slide 8: Go-to-Market
    slides_data.append({
        'type': 'content',
        'title': 'Go-to-Market: Salesforce SMB Entry Wedge',
        'content': [
            {'text': 'Phase 1: Salesforce Service Cloud SMBs (Yr 1-2)', 'level': 0},
            {'text': 'AppExchange freemium, $375/agent/mo, <1 day setup', 'level': 1},
            {'text': 'Target: 200 customers, 2K agents, $8.4M ARR Year 1', 'level': 1},
            {'text': 'Salesforce sellers as advocates (unstick deals)', 'level': 1},
            {'text': 'Phase 2: Mid-Market + Sales Cloud (Yr 2-3)', 'level': 0},
            {'text': 'Expand to mid-market Service Cloud', 'level': 1},
            {'text': 'Cross-sell Sales Cloud automation', 'level': 1},
            {'text': 'Target: 1,400 customers, $71M ARR Year 3', 'level': 1},
            {'text': 'Phase 3: Beyond Salesforce (Yr 3-5)', 'level': 0},
            {'text': 'ServiceNow, Zendesk, Microsoft Dynamics', 'level': 1},
            {'text': 'Target: $212M ARR Year 5 across platforms', 'level': 1},
        ]
    })

    # Slide 9: Product Roadmap Overview
    slides_data.append({
        'type': 'content',
        'title': 'Product Roadmap: 18-Month Path',
        'content': [
            {'text': 'Q4 2025: Foundation', 'level': 0},
            {'text': 'Desktop app, edge LLM, basic peer benchmarking', 'level': 1},
            {'text': 'Target: 5K free users, 50 design partners', 'level': 1},
            {'text': 'Q1 2026: Monetization', 'level': 0},
            {'text': 'AI roadmaps, agent authoring, token billing', 'level': 1},
            {'text': 'Target: 15K free, 100 paying, $250K MRR', 'level': 1},
            {'text': 'Q2 2026: Scale', 'level': 0},
            {'text': 'Performance monitoring, continuous improvement', 'level': 1},
            {'text': 'Target: 40K free, 600 paying, $1M MRR', 'level': 1},
            {'text': 'Q3-Q4 2026: Enterprise Platform', 'level': 0},
            {'text': 'Governance console, SOC 2, enterprise features', 'level': 1},
            {'text': 'Target: 100K free, 3K paying, $6M ARR', 'level': 1},
        ]
    })

    # Slide 10: Product Roadmap Matrix - 2026 by Quarter
    slides_data.append({
        'type': 'table',
        'title': 'Product Roadmap by Quarter: 2026',
        'headers': ["Phase", "Q1 2026", "Q2 2026", "Q3 2026", "Q4 2026"],
        'rows': [
            [
                "MINE\n(Discovery)",
                "• Enhanced workflow capture\n• Multi-user synthesis\n• Team-level process mapping\n• Timeline visualization",
                "• Department-level discovery\n• Cross-functional mapping\n• Historical trend analysis\n• Process change detection",
                "• Enterprise-wide discovery\n• Multi-system tracking\n• Compliance workflow mapping\n• Advanced pattern recognition",
                "• Real-time observation\n• Predictive process analysis\n• Industry-specific templates\n• API for external data"
            ],
            [
                "DESIGN\n(Intelligence)",
                "• AI transformation roadmaps\n• Effort vs. impact matrices\n• Success pattern matching\n• Basic peer benchmarking",
                "• Advanced benchmarking (20+ industries)\n• Custom benchmark reports\n• ROI scenario modeling\n• What-if analysis",
                "• Competitive intelligence\n• Cross-functional optimization\n• Strategic transformation planning\n• Executive dashboards",
                "• Predictive outcomes\n• Best practice library\n• Custom KPI frameworks\n• Strategic advisory insights"
            ],
            [
                "DELIVER\n(Execution)",
                "• Autonomous agent authoring\n• Shadow mode validation\n• Agent templates library\n• Token-based deployment",
                "• Agent performance optimization\n• Multi-agent orchestration\n• Version control for agents\n• A/B testing framework",
                "• Enterprise agent governance\n• Custom integration framework\n• Agent marketplace foundation\n• Rollback & recovery",
                "• Advanced agent analytics\n• Self-healing agents\n• Partner agent ecosystem\n• Low-code agent builder"
            ],
            [
                "MANAGE\n(Continuous Improvement)",
                "• Basic performance monitoring\n• Simple alerting\n• Weekly reports\n• Usage analytics",
                "• Real-time monitoring dashboards\n• Service assurance automation\n• Continuous improvement recommendations\n• Automated executive reporting",
                "• Enterprise governance console\n• Multi-team coordination\n• SLA management\n• Advanced anomaly detection",
                "• Predictive service assurance\n• Automated optimization cycles\n• Strategic performance analytics\n• Board-level reporting"
            ]
        ]
    })

    # Slide 11: Financial Projections
    slides_data.append({
        'type': 'table',
        'title': 'Financial Projections: Path to $300M ARR',
        'headers': ["Year", "SMB", "Mid-Market", "Enterprise", "Total ARR"],
        'rows': [
            ["1", "$1.2M", "$0.5M", "$0.3M", "$2M"],
            ["2", "$4.5M", "$3.0M", "$2.5M", "$10M"],
            ["3", "$12M", "$14M", "$14M", "$40M"],
            ["4", "$26M", "$38M", "$56M", "$120M"],
            ["5", "$50M", "$80M", "$170M", "$300M"]
        ]
    })

    # Slide 12: Team & Capital
    slides_data.append({
        'type': 'content',
        'title': 'Team & Capital Requirements',
        'content': [
            {'text': 'Founding Team:', 'level': 0},
            {'text': 'CEO: Vision, fundraising, partnerships', 'level': 1},
            {'text': 'Product Owner: Platform architecture, AI strategy', 'level': 1},
            {'text': 'Technical Lead: Edge AI, privacy systems', 'level': 1},
            {'text': 'Phase 1 Hires (Mo 1-6): 8-12 people', 'level': 0},
            {'text': 'ML Engineer, Growth Lead, Operations Lead, Engineers', 'level': 1},
            {'text': 'Capital Requirements: $2-3M Seed Extension', 'level': 0},
            {'text': 'Team (60%): $1.4-1.8M', 'level': 1},
            {'text': 'Infrastructure (20%): $500-750K', 'level': 1},
            {'text': 'Go-to-Market (15%): $375-600K', 'level': 1},
            {'text': 'R&D (5%): $125-300K', 'level': 1},
        ]
    })

    # Slide 13: The Ask
    slides_data.append({
        'type': 'content',
        'title': 'The Ask: $2-3M Seed Extension',
        'content': [
            {'text': 'Milestones This Enables:', 'level': 0},
            {'text': 'Mo 3: 5K free users, 50 design partners, PMF', 'level': 1},
            {'text': 'Mo 6: $250K MRR, 100 paying, 30% conversion', 'level': 1},
            {'text': 'Mo 12: $1.5M MRR, 1.5K customers', 'level': 1},
            {'text': 'Mo 18: $6M ARR, 3K+ customers, Series A ready', 'level': 1},
            {'text': 'Why This Wins:', 'level': 0},
            {'text': 'Low-end disruption: 150K SMBs incumbents ignore', 'level': 1},
            {'text': 'Network effects: Data moat compounds daily', 'level': 1},
            {'text': 'Capital efficient: 229:1 LTV:CAC, 85% margins', 'level': 1},
            {'text': 'Perfect timing: Agent economy + SMB transformation', 'level': 1},
            {'text': 'We\'re democratizing transformation. Building the category leader.', 'level': 0},
        ]
    })

    return slides_data

def archive_existing_file(file_path, archive_dir):
    """Archive existing file with timestamp to avoid collisions."""
    if os.path.exists(file_path):
        # Create archive directory if it doesn't exist
        os.makedirs(archive_dir, exist_ok=True)

        # Generate timestamped filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        base_name = os.path.basename(file_path)
        name_without_ext = os.path.splitext(base_name)[0]
        ext = os.path.splitext(base_name)[1]

        archive_name = f"{name_without_ext}_{timestamp}{ext}"
        archive_path = os.path.join(archive_dir, archive_name)

        # Copy file to archive
        shutil.copy2(file_path, archive_path)
        print(f"Archived existing file to: {archive_name}")
        return archive_path
    return None

def main():
    # Determine project root (parent of scripts directory or current if we're already at root)
    current_dir = os.path.dirname(os.path.abspath(__file__))
    if os.path.basename(current_dir) == 'scripts':
        project_root = os.path.dirname(current_dir)
    else:
        project_root = current_dir

    output_path = os.path.join(project_root, "outputdocs", "EVOLVE_InvestorDeck.pptx")
    archive_dir = os.path.join(project_root, "outputdocs", "archive")

    # Archive existing file before making changes
    print("Checking for existing presentation...")
    archive_existing_file(output_path, archive_dir)

    # Load existing presentation to preserve title slide and styling
    print("Loading existing presentation to preserve title slide and styling...")
    try:
        prs = Presentation(output_path)
        print(f"Loaded existing presentation with {len(prs.slides)} slides and {len(prs.slide_layouts)} layouts")

        # Extract theme from existing presentation
        theme = extract_styling_from_example(output_path)

        # Keep only the title slide (slide 0), remove all others
        print("Preserving title slide, removing other slides...")
        while len(prs.slides) > 1:
            rId = prs.slides._sldIdLst[1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[1]

        print(f"Title slide preserved. Ready to add new content slides.")

    except FileNotFoundError:
        # If file doesn't exist, use board deck template as fallback
        print("No existing presentation found. Using board deck template...")
        template_path = os.path.join(project_root, "styles", "Synaptic-BODMeeting-2025-10-23-MgmtSession.pptx")
        prs = Presentation(template_path)
        theme = extract_styling_from_example(template_path)

        # Clear all slides from template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

    print("Parsing investor brief content...")
    slides_data = parse_investor_brief_content()

    # Skip the first slide data (title) since we're preserving the existing title slide
    print(f"Creating {len(slides_data) - 1} new content slides (keeping existing title slide)...")
    print(f"Total deck will have {len(slides_data)} slides (1 preserved title + {len(slides_data) - 1} new slides)")
    for i, slide_data in enumerate(slides_data[1:], 2):  # Start from slide 2, skip title slide
        print(f"Creating slide {i}: {slide_data.get('title', 'Content slide')}")

        if slide_data['type'] == 'content':
            slide = create_content_slide(prs, slide_data['title'], slide_data['content'], theme)

        elif slide_data['type'] == 'table':
            slide = create_table_slide(prs, slide_data['title'], slide_data['headers'], slide_data['rows'], theme)

    print(f"\nSaving presentation to {output_path}...")
    prs.save(output_path)
    print("Done! Investor deck created successfully.")

if __name__ == "__main__":
    main()
